#!/usr/bin/env python3
"""
importar_plano.py — extrai as tabelas dos PPTX de planejamento de safra (modelo do agrônomo)
para um JSON bruto. NÃO faz de-para de unidades: os nomes saem exatamente como estão nos slides.
O de-para é feito depois, contra a tabela unidade_alias (sistema = 'plano').

Uso:  python3 importar_plano.py pasta_com_pptx/ saida.json
Requer: python-pptx  (pip install python-pptx)

Lê: estimativa (talhão, área, poda, renovação), fito mês a mês, correção de solo,
resumo de adubos, calendários mensais por unidade (kg), Gantt (pela COR das células).
Testado em 03/09/2026 nos 7 decks da safra 2026/27.
"""
import sys, re, json, glob, os
from pptx import Presentation

NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
INSUMOS = ['ureia', 'nitrato', 'sulfato_amonio', 'kcl', 'phusion', 'sulfato_mn', 'acido_borico', 'sulfato_zn']
MESES = {'setembro': 9, 'outubro': 10, 'novembro': 11, 'dezembro': 12, 'janeiro': 1, 'fevereiro': 2,
         'março': 3, 'marco': 3, 'abril': 4, 'maio': 5, 'junho': 6, 'julho': 7}
MES_ABREV = {'ago': 8, 'set': 9, 'out': 10, 'nov': 11, 'dez': 12, 'jan': 1, 'fev': 2, 'mar': 3, 'abr': 4, 'mai': 5, 'jun': 6, 'jul': 7}
FITO_MESES = ['agosto', 'setembro', 'outubro', 'novembro', 'dezembro', 'janeiro', 'fevereiro', 'março', 'abril', 'maio']
FITO_NUM = [8, 9, 10, 11, 12, 1, 2, 3, 4, 5]


def num(s):
    s = (s or '').strip().replace('.', '')
    m = re.match(r'^-?\d+(,\d+)?', s)
    return float(m.group(0).replace(',', '.')) if m else None


def cell_fill(c):
    tcPr = c._tc.find('a:tcPr', NS)
    if tcPr is None:
        return None
    sf = tcPr.find('a:solidFill', NS)
    return None if sf is None else sf[0].get('val')


def slide_title(slide):
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip() and not shp.has_table:
            return shp.text_frame.text.strip().replace('\n', ' ').replace('\x0b', ' ')
    return ''


def ler_deck(path):
    prs = Presentation(path)
    out = {'arquivo': os.path.basename(path), 'capa': '', 'estimativa': [], 'fito': [], 'calagem': [],
           'resumo': [], 'calendarios': [], 'gantt': {}, 'avisos': []}
    for idx, slide in enumerate(prs.slides, 1):
        title = slide_title(slide)
        if idx == 1:
            out['capa'] = ' | '.join(s.text_frame.text.strip().replace('\x0b', ' ') for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
        for shp in slide.shapes:
            if not shp.has_table:
                continue
            tbl = shp.table
            rows = [[c.text.strip() for c in r.cells] for r in tbl.rows]
            hdr = [h.lower() for h in rows[0]]
            first = hdr[0] if hdr else ''
            # ---- calendário mensal por unidade
            if any(r[0].strip().lower() in MESES for r in rows):
                unidade = re.sub(r'^calend[áa]rio\s+aduba[çc][ãa]o\s*', '', title, flags=re.I).strip() or title
                cal = {'slide': idx, 'unidade_nome': unidade, 'meses': []}
                for r in rows[1:]:
                    m = r[0].strip().lower()
                    if m not in MESES:
                        continue
                    vals = r[1:9] + [''] * (8 - len(r[1:9]))
                    for i, v in enumerate(vals):
                        if v.strip():
                            n = num(v)
                            if n is None:
                                out['avisos'].append(f'slide {idx} {unidade} {m}/{INSUMOS[i]}: valor não numérico "{v}"')
                                continue
                            item = {'mes': MESES[m], 'insumo': INSUMOS[i], 'kg': n}
                            if not re.fullmatch(r'\s*[\d\.]+\s*', v):
                                item['obs'] = v
                            cal['meses'].append(item)
                if len(hdr) != 9:
                    out['avisos'].append(f'slide {idx} {unidade}: calendário com {len(hdr)} colunas (esperado 9)')
                out['calendarios'].append(cal)
            # ---- fito mês a mês
            elif first == 'agosto':
                if len(rows) >= 3:
                    for i, m in enumerate(rows[0]):
                        out['fito'].append({'mes': FITO_NUM[i] if i < 10 else None, 'mes_nome': m,
                                            'fase_alvos': rows[1][i] if i < len(rows[1]) else '',
                                            'produtos': rows[2][i] if i < len(rows[2]) else ''})
            # ---- Gantt (ano agrícola)
            elif first.startswith('ano agr'):
                meses = [h.strip().lower()[:3] for h in rows[0][1:]]
                for r in list(tbl.rows)[1:]:
                    cells = list(r.cells)
                    nome = cells[0].text.strip()
                    if not nome or nome.lower().startswith('eng'):
                        continue
                    marcados = [MES_ABREV[meses[i]] for i, c in enumerate(cells[1:]) if i < len(meses) and cell_fill(c) not in (None, 'window', 'bg1', 'FFFFFF')]
                    out['gantt'][nome] = marcados
            # ---- resumo de adubos
            elif first == 'fazenda':
                for r in rows[1:]:
                    if r[0].strip():
                        vals = r[1:9] + [''] * (8 - len(r[1:9]))
                        out['resumo'].append({'rotulo': r[0].strip(), **{INSUMOS[i]: (num(v) or 0) for i, v in enumerate(vals)}})
            # ---- correção de solo
            elif 'talh' in first and any('calc' in h for h in hdr):
                for r in rows[1:]:
                    if r[0].strip():
                        out['calagem'].append({'unidade_nome': r[0].strip(), 't_ha': num(r[1]) if len(r) > 1 else None,
                                               't_total': num(r[2]) if len(r) > 2 else None, 'total_linha': r[0].strip().lower().startswith(('total', 'resumo'))})
            # ---- estimativa de produção
            elif 'talh' in first and any('rea' in h for h in hdr):
                for r in rows[1:]:
                    if r[0].strip():
                        out['estimativa'].append({'unidade_nome': r[0].strip(), 'area_ha': num(r[1]) if len(r) > 1 else None,
                                                  'media_producao': r[2] if len(r) > 2 else '', 'producao_total': r[3] if len(r) > 3 else '',
                                                  'poda_ha': num(r[4]) if len(r) > 4 else None, 'renovacao_ha': num(r[5]) if len(r) > 5 else None,
                                                  'total_linha': r[0].strip().lower() == 'total'})
    # ---- auditoria mínima
    soma = {i: sum(m['kg'] for c in out['calendarios'] for m in c['meses'] if m['insumo'] == i) for i in INSUMOS}
    out['soma_calendarios'] = soma
    if len(out['resumo']) == 1:
        dif = {i: out['resumo'][0][i] - soma[i] for i in INSUMOS if out['resumo'][0][i] != soma[i]}
        if dif:
            out['avisos'].append(f'resumo − soma dos calendários: {dif}')
    areas = [e['area_ha'] for e in out['estimativa'] if not e['total_linha'] and e['area_ha']]
    tot = [e['area_ha'] for e in out['estimativa'] if e['total_linha']]
    if areas and tot and tot[0] and abs(sum(areas) - tot[0]) > 0.5:
        out['avisos'].append(f'soma das áreas {sum(areas)} ≠ total do slide {tot[0]}')
    for c in out['calagem']:
        if not c['total_linha'] and c['t_ha'] and c['t_total']:
            c['area_implicita_ha'] = round(c['t_total'] / c['t_ha'], 1)
    return out


def main():
    if len(sys.argv) < 3:
        print(__doc__); sys.exit(1)
    pasta, saida = sys.argv[1], sys.argv[2]
    decks = [ler_deck(p) for p in sorted(glob.glob(os.path.join(pasta, '*.pptx')))]
    json.dump({'decks': decks}, open(saida, 'w'), ensure_ascii=False, indent=1)
    for d in decks:
        print(f"{d['arquivo']}: {len(d['calendarios'])} calendários, {len(d['calagem'])} linhas de calagem, "
              f"{len(d['estimativa'])} de estimativa, {len(d['gantt'])} linhas de Gantt, {len(d['avisos'])} avisos")
        for a in d['avisos']:
            print('   ⚠', a)


if __name__ == '__main__':
    main()
